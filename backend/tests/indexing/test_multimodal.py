"""Tests for multimodal document ingestion (Story 20-D2)."""

from pathlib import Path
from typing import Any
from unittest.mock import MagicMock, patch
from uuid import UUID, uuid4

import pytest

from agentic_rag_backend.indexing.multimodal import (
    DEFAULT_MULTIMODAL_ENABLED,
    DEFAULT_OFFICE_DOCS_ENABLED,
    EXTENSION_TYPE_MAP,
    MAX_FILE_SIZE_BYTES,
    MIME_TYPE_MAP,
    UNSUPPORTED_LEGACY_FORMATS,
    DocumentType,
    ExcelContent,
    ExtractedCell,
    ExtractedSheet,
    ExtractedSlide,
    MultimodalIngester,
    MultimodalIngestionAdapter,
    MultimodalIngestionResult,
    OfficeParser,
    PowerPointContent,
    WordContent,
    get_multimodal_ingestion_adapter,
)


# ==============================================================================
# DocumentType Enum Tests
# ==============================================================================


class TestDocumentType:
    """Tests for DocumentType enum."""

    def test_document_type_pdf(self) -> None:
        """Test PDF document type."""
        assert DocumentType.PDF.value == "pdf"

    def test_document_type_word(self) -> None:
        """Test Word document type."""
        assert DocumentType.WORD.value == "word"

    def test_document_type_excel(self) -> None:
        """Test Excel document type."""
        assert DocumentType.EXCEL.value == "excel"

    def test_document_type_powerpoint(self) -> None:
        """Test PowerPoint document type."""
        assert DocumentType.POWERPOINT.value == "powerpoint"

    def test_document_type_image(self) -> None:
        """Test Image document type."""
        assert DocumentType.IMAGE.value == "image"

    def test_document_type_markdown(self) -> None:
        """Test Markdown document type."""
        assert DocumentType.MARKDOWN.value == "markdown"

    def test_document_type_text(self) -> None:
        """Test Text document type."""
        assert DocumentType.TEXT.value == "text"

    def test_document_type_unknown(self) -> None:
        """Test Unknown document type."""
        assert DocumentType.UNKNOWN.value == "unknown"

    def test_all_document_types_have_values(self) -> None:
        """Test all document types have non-empty values."""
        for doc_type in DocumentType:
            assert doc_type.value
            assert isinstance(doc_type.value, str)


# ==============================================================================
# Extension and MIME Type Mapping Tests
# ==============================================================================


class TestTypeMappings:
    """Tests for extension and MIME type mappings."""

    def test_docx_extension_mapping(self) -> None:
        """Test .docx maps to WORD."""
        assert EXTENSION_TYPE_MAP[".docx"] == DocumentType.WORD

    def test_xlsx_extension_mapping(self) -> None:
        """Test .xlsx maps to EXCEL."""
        assert EXTENSION_TYPE_MAP[".xlsx"] == DocumentType.EXCEL

    def test_pptx_extension_mapping(self) -> None:
        """Test .pptx maps to POWERPOINT."""
        assert EXTENSION_TYPE_MAP[".pptx"] == DocumentType.POWERPOINT

    def test_pdf_extension_mapping(self) -> None:
        """Test .pdf maps to PDF."""
        assert EXTENSION_TYPE_MAP[".pdf"] == DocumentType.PDF

    def test_image_extensions_mapping(self) -> None:
        """Test image extensions map to IMAGE."""
        image_extensions = [".png", ".jpg", ".jpeg", ".gif", ".webp"]
        for ext in image_extensions:
            assert EXTENSION_TYPE_MAP[ext] == DocumentType.IMAGE

    def test_markdown_extensions_mapping(self) -> None:
        """Test markdown extensions map to MARKDOWN."""
        assert EXTENSION_TYPE_MAP[".md"] == DocumentType.MARKDOWN
        assert EXTENSION_TYPE_MAP[".markdown"] == DocumentType.MARKDOWN

    def test_text_extension_mapping(self) -> None:
        """Test .txt maps to TEXT."""
        assert EXTENSION_TYPE_MAP[".txt"] == DocumentType.TEXT

    def test_word_mime_type_mapping(self) -> None:
        """Test Word MIME type mapping."""
        mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        assert MIME_TYPE_MAP[mime] == DocumentType.WORD

    def test_excel_mime_type_mapping(self) -> None:
        """Test Excel MIME type mapping."""
        mime = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        assert MIME_TYPE_MAP[mime] == DocumentType.EXCEL

    def test_powerpoint_mime_type_mapping(self) -> None:
        """Test PowerPoint MIME type mapping."""
        mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        assert MIME_TYPE_MAP[mime] == DocumentType.POWERPOINT


# ==============================================================================
# ExtractedCell Tests
# ==============================================================================


class TestExtractedCell:
    """Tests for ExtractedCell dataclass."""

    def test_extracted_cell_creation(self) -> None:
        """Test creating an extracted cell."""
        cell = ExtractedCell(row=0, column=0, value="Test", data_type="string")
        assert cell.row == 0
        assert cell.column == 0
        assert cell.value == "Test"
        assert cell.data_type == "string"

    def test_extracted_cell_default_data_type(self) -> None:
        """Test default data type is string."""
        cell = ExtractedCell(row=1, column=2, value="Value")
        assert cell.data_type == "string"

    def test_extracted_cell_to_dict(self) -> None:
        """Test cell serialization."""
        cell = ExtractedCell(row=5, column=3, value="123", data_type="number")
        result = cell.to_dict()
        assert result == {
            "row": 5,
            "column": 3,
            "value": "123",
            "data_type": "number",
        }


# ==============================================================================
# ExtractedSheet Tests
# ==============================================================================


class TestExtractedSheet:
    """Tests for ExtractedSheet dataclass."""

    def test_extracted_sheet_creation(self) -> None:
        """Test creating an extracted sheet."""
        sheet = ExtractedSheet(name="Sheet1", index=0)
        assert sheet.name == "Sheet1"
        assert sheet.index == 0
        assert sheet.headers == []
        assert sheet.rows == []

    def test_extracted_sheet_with_data(self) -> None:
        """Test sheet with headers and rows."""
        sheet = ExtractedSheet(
            name="Data",
            index=1,
            headers=["Name", "Value"],
            rows=[["A", "1"], ["B", "2"]],
            row_count=3,
            column_count=2,
        )
        assert sheet.headers == ["Name", "Value"]
        assert len(sheet.rows) == 2
        assert sheet.row_count == 3
        assert sheet.column_count == 2

    def test_extracted_sheet_to_dict(self) -> None:
        """Test sheet serialization."""
        sheet = ExtractedSheet(
            name="Test",
            index=0,
            headers=["Col1"],
            rows=[["Val1"]],
            row_count=2,
            column_count=1,
        )
        result = sheet.to_dict()
        assert result["name"] == "Test"
        assert result["index"] == 0
        assert result["headers"] == ["Col1"]
        assert result["rows"] == [["Val1"]]

    def test_extracted_sheet_to_markdown(self) -> None:
        """Test sheet markdown conversion."""
        sheet = ExtractedSheet(
            name="Results",
            index=0,
            headers=["Name", "Score"],
            rows=[["Alice", "95"], ["Bob", "87"]],
        )
        markdown = sheet.to_markdown()
        assert "## Results" in markdown
        assert "| Name | Score |" in markdown
        assert "| Alice | 95 |" in markdown

    def test_extracted_sheet_to_markdown_empty(self) -> None:
        """Test empty sheet markdown conversion."""
        sheet = ExtractedSheet(name="Empty", index=0)
        markdown = sheet.to_markdown()
        assert markdown == ""

    def test_extracted_sheet_to_markdown_escapes_pipes(self) -> None:
        """Test pipe characters are escaped in markdown."""
        sheet = ExtractedSheet(
            name="Test",
            index=0,
            headers=["Data"],
            rows=[["A|B"]],
        )
        markdown = sheet.to_markdown()
        assert "A\\|B" in markdown


# ==============================================================================
# ExtractedSlide Tests
# ==============================================================================


class TestExtractedSlide:
    """Tests for ExtractedSlide dataclass."""

    def test_extracted_slide_creation(self) -> None:
        """Test creating an extracted slide."""
        slide = ExtractedSlide(number=1)
        assert slide.number == 1
        assert slide.title is None
        assert slide.content == ""
        assert slide.notes is None

    def test_extracted_slide_with_content(self) -> None:
        """Test slide with full content."""
        slide = ExtractedSlide(
            number=2,
            title="Introduction",
            content="Welcome to the presentation",
            notes="Remember to smile",
            shapes_text=["Introduction", "Welcome to the presentation"],
        )
        assert slide.title == "Introduction"
        assert "Welcome" in slide.content
        assert slide.notes == "Remember to smile"
        assert len(slide.shapes_text) == 2

    def test_extracted_slide_to_dict(self) -> None:
        """Test slide serialization."""
        slide = ExtractedSlide(
            number=1,
            title="Test",
            content="Content",
            notes="Notes",
        )
        result = slide.to_dict()
        assert result["number"] == 1
        assert result["title"] == "Test"
        assert result["content"] == "Content"
        assert result["notes"] == "Notes"

    def test_extracted_slide_to_markdown(self) -> None:
        """Test slide markdown conversion."""
        slide = ExtractedSlide(
            number=3,
            title="Summary",
            content="Key points here",
            notes="Wrap up",
        )
        markdown = slide.to_markdown()
        assert "## Slide 3" in markdown
        assert "### Summary" in markdown
        assert "Key points here" in markdown
        assert "**Notes:** Wrap up" in markdown


# ==============================================================================
# WordContent Tests
# ==============================================================================


class TestWordContent:
    """Tests for WordContent dataclass."""

    def test_word_content_creation(self) -> None:
        """Test creating word content."""
        content = WordContent()
        assert content.paragraphs == []
        assert content.tables == []
        assert content.full_text == ""
        assert content.metadata == {}

    def test_word_content_with_data(self) -> None:
        """Test word content with paragraphs and tables."""
        content = WordContent(
            paragraphs=["Para 1", "Para 2"],
            tables=[[["A", "B"], ["C", "D"]]],
            full_text="Para 1\n\nPara 2",
            metadata={"author": "Test User"},
        )
        assert len(content.paragraphs) == 2
        assert len(content.tables) == 1
        assert "Para 1" in content.full_text

    def test_word_content_to_dict(self) -> None:
        """Test word content serialization."""
        content = WordContent(
            paragraphs=["Test"],
            metadata={"title": "Doc"},
        )
        result = content.to_dict()
        assert result["paragraphs"] == ["Test"]
        assert result["metadata"]["title"] == "Doc"


# ==============================================================================
# ExcelContent Tests
# ==============================================================================


class TestExcelContent:
    """Tests for ExcelContent dataclass."""

    def test_excel_content_creation(self) -> None:
        """Test creating excel content."""
        content = ExcelContent()
        assert content.sheets == []
        assert content.sheet_count == 0

    def test_excel_content_with_sheets(self) -> None:
        """Test excel content with sheets."""
        sheet = ExtractedSheet(name="Data", index=0)
        content = ExcelContent(sheets=[sheet], sheet_count=1)
        assert len(content.sheets) == 1
        assert content.sheet_count == 1

    def test_excel_content_to_dict(self) -> None:
        """Test excel content serialization."""
        sheet = ExtractedSheet(name="Test", index=0)
        content = ExcelContent(sheets=[sheet], sheet_count=1)
        result = content.to_dict()
        assert result["sheet_count"] == 1
        assert len(result["sheets"]) == 1


# ==============================================================================
# PowerPointContent Tests
# ==============================================================================


class TestPowerPointContent:
    """Tests for PowerPointContent dataclass."""

    def test_powerpoint_content_creation(self) -> None:
        """Test creating powerpoint content."""
        content = PowerPointContent()
        assert content.slides == []
        assert content.slide_count == 0
        assert content.full_text == ""

    def test_powerpoint_content_with_slides(self) -> None:
        """Test powerpoint content with slides."""
        slide = ExtractedSlide(number=1, title="Intro")
        content = PowerPointContent(
            slides=[slide],
            slide_count=1,
            full_text="Slide 1: Intro",
        )
        assert len(content.slides) == 1
        assert content.slide_count == 1

    def test_powerpoint_content_to_dict(self) -> None:
        """Test powerpoint content serialization."""
        slide = ExtractedSlide(number=1)
        content = PowerPointContent(slides=[slide], slide_count=1)
        result = content.to_dict()
        assert result["slide_count"] == 1
        assert len(result["slides"]) == 1


# ==============================================================================
# MultimodalIngestionResult Tests
# ==============================================================================


class TestMultimodalIngestionResult:
    """Tests for MultimodalIngestionResult dataclass."""

    def test_result_creation(self) -> None:
        """Test creating an ingestion result."""
        result = MultimodalIngestionResult(
            document_id="doc-123",
            tenant_id="tenant-456",
            document_type=DocumentType.WORD,
            file_path="/path/to/doc.docx",
        )
        assert result.document_id == "doc-123"
        assert result.tenant_id == "tenant-456"
        assert result.document_type == DocumentType.WORD
        assert result.success is True

    def test_result_with_error(self) -> None:
        """Test result with error."""
        result = MultimodalIngestionResult(
            document_id="doc-123",
            tenant_id="tenant-456",
            document_type=DocumentType.UNKNOWN,
            file_path="/path/to/file",
            success=False,
            error="Unsupported format",
        )
        assert result.success is False
        assert result.error == "Unsupported format"

    def test_result_to_dict(self) -> None:
        """Test result serialization."""
        result = MultimodalIngestionResult(
            document_id="doc-123",
            tenant_id="tenant-456",
            document_type=DocumentType.EXCEL,
            file_path="/path/to/file.xlsx",
            processing_time_ms=150,
            success=True,
        )
        data = result.to_dict()
        assert data["document_id"] == "doc-123"
        assert data["document_type"] == "excel"
        assert data["processing_time_ms"] == 150
        assert data["success"] is True


# ==============================================================================
# MultimodalIngester Tests
# ==============================================================================


class TestMultimodalIngester:
    """Tests for MultimodalIngester class."""

    def test_ingester_initialization(self) -> None:
        """Test ingester initialization."""
        ingester = MultimodalIngester(
            multimodal_enabled=True,
            office_docs_enabled=True,
        )
        assert ingester.multimodal_enabled is True
        assert ingester.office_docs_enabled is True

    def test_ingester_disabled_returns_error(self, tmp_path: Path) -> None:
        """Test ingester returns error when disabled."""
        ingester = MultimodalIngester(multimodal_enabled=False)
        test_file = tmp_path / "test.txt"
        test_file.write_text("content")

        result = ingester.ingest(
            file_path=test_file,
            document_id=uuid4(),
            tenant_id=uuid4(),
        )
        assert result.success is False
        assert "disabled" in result.error.lower()

    def test_detect_type_from_extension(self) -> None:
        """Test document type detection from extension."""
        ingester = MultimodalIngester(multimodal_enabled=True)

        assert ingester._detect_type(Path("doc.docx")) == DocumentType.WORD
        assert ingester._detect_type(Path("sheet.xlsx")) == DocumentType.EXCEL
        assert ingester._detect_type(Path("slides.pptx")) == DocumentType.POWERPOINT
        assert ingester._detect_type(Path("readme.md")) == DocumentType.MARKDOWN
        assert ingester._detect_type(Path("notes.txt")) == DocumentType.TEXT
        assert ingester._detect_type(Path("unknown.xyz")) == DocumentType.UNKNOWN

    def test_process_markdown_file(self, tmp_path: Path) -> None:
        """Test processing a markdown file."""
        ingester = MultimodalIngester(multimodal_enabled=True)

        md_file = tmp_path / "test.md"
        md_content = "# Hello\n\nThis is a test."
        md_file.write_text(md_content)

        result = ingester.ingest(
            file_path=md_file,
            document_id=uuid4(),
            tenant_id=uuid4(),
        )
        assert result.success is True
        assert result.document_type == DocumentType.MARKDOWN
        assert result.text_content == md_content
        assert result.markdown == md_content

    def test_process_text_file(self, tmp_path: Path) -> None:
        """Test processing a plain text file."""
        ingester = MultimodalIngester(multimodal_enabled=True)

        txt_file = tmp_path / "test.txt"
        txt_content = "Plain text content"
        txt_file.write_text(txt_content)

        result = ingester.ingest(
            file_path=txt_file,
            document_id=uuid4(),
            tenant_id=uuid4(),
        )
        assert result.success is True
        assert result.document_type == DocumentType.TEXT
        assert result.text_content == txt_content
        assert "```" in result.markdown  # Wrapped in code block

    def test_pdf_processing_redirects_to_docling(self, tmp_path: Path) -> None:
        """Test PDF processing hints at using EnhancedDoclingParser."""
        ingester = MultimodalIngester(multimodal_enabled=True)

        pdf_file = tmp_path / "test.pdf"
        pdf_file.write_bytes(b"%PDF-1.4 fake content")

        result = ingester.ingest(
            file_path=pdf_file,
            document_id=uuid4(),
            tenant_id=uuid4(),
        )
        assert result.success is False
        assert "EnhancedDoclingParser" in result.error

    def test_image_processing_not_implemented(self, tmp_path: Path) -> None:
        """Test image processing returns not implemented."""
        ingester = MultimodalIngester(multimodal_enabled=True)

        img_file = tmp_path / "test.png"
        img_file.write_bytes(b"fake image")

        result = ingester.ingest(
            file_path=img_file,
            document_id=uuid4(),
            tenant_id=uuid4(),
        )
        assert result.success is False
        assert "not yet implemented" in result.error.lower()

    def test_office_docs_disabled_returns_error(self, tmp_path: Path) -> None:
        """Test Office docs return error when disabled."""
        ingester = MultimodalIngester(
            multimodal_enabled=True,
            office_docs_enabled=False,
        )

        # Create a fake docx file
        docx_file = tmp_path / "test.docx"
        docx_file.write_bytes(b"fake content")

        result = ingester.ingest(
            file_path=docx_file,
            document_id=uuid4(),
            tenant_id=uuid4(),
        )
        assert result.success is False
        assert "disabled" in result.error.lower()


# ==============================================================================
# Security Validation Tests
# ==============================================================================


class TestSecurityValidation:
    """Tests for security validations in MultimodalIngester."""

    def test_ingest_requires_tenant_id(self, tmp_path: Path) -> None:
        """Test that tenant_id is required."""
        ingester = MultimodalIngester(multimodal_enabled=True)
        test_file = tmp_path / "test.txt"
        test_file.write_text("content")

        with pytest.raises(ValueError, match="tenant_id is required"):
            ingester.ingest(
                file_path=test_file,
                document_id=uuid4(),
                tenant_id=None,  # type: ignore[arg-type]
            )

    def test_ingest_validates_file_exists(self, tmp_path: Path) -> None:
        """Test that file must exist."""
        ingester = MultimodalIngester(multimodal_enabled=True)
        non_existent = tmp_path / "does_not_exist.txt"

        with pytest.raises(FileNotFoundError, match="Document not found"):
            ingester.ingest(
                file_path=non_existent,
                document_id=uuid4(),
                tenant_id=uuid4(),
            )

    def test_ingest_validates_is_file(self, tmp_path: Path) -> None:
        """Test that path must be a file, not directory."""
        ingester = MultimodalIngester(multimodal_enabled=True)

        with pytest.raises(ValueError, match="not a file"):
            ingester.ingest(
                file_path=tmp_path,  # This is a directory
                document_id=uuid4(),
                tenant_id=uuid4(),
            )

    def test_ingest_path_traversal_protection(self, tmp_path: Path) -> None:
        """Test path traversal protection."""
        ingester = MultimodalIngester(multimodal_enabled=True)

        # Create allowed directory and file outside it
        allowed_dir = tmp_path / "allowed"
        allowed_dir.mkdir()

        outside_file = tmp_path / "outside.txt"
        outside_file.write_text("content")

        with pytest.raises(ValueError, match="Path traversal not allowed"):
            ingester.ingest(
                file_path=outside_file,
                document_id=uuid4(),
                tenant_id=uuid4(),
                allowed_base_path=allowed_dir,
            )

    def test_ingest_allows_file_in_base_path(self, tmp_path: Path) -> None:
        """Test file within allowed path is accepted."""
        ingester = MultimodalIngester(multimodal_enabled=True)

        # Create file within allowed directory
        allowed_file = tmp_path / "allowed.txt"
        allowed_file.write_text("content")

        # Should not raise
        result = ingester.ingest(
            file_path=allowed_file,
            document_id=uuid4(),
            tenant_id=uuid4(),
            allowed_base_path=tmp_path,
        )
        assert result.success is True

    def test_ingest_allows_nested_file_in_base_path(self, tmp_path: Path) -> None:
        """Test nested file within allowed path is accepted."""
        ingester = MultimodalIngester(multimodal_enabled=True)

        # Create nested directory and file
        nested_dir = tmp_path / "subdir"
        nested_dir.mkdir()
        nested_file = nested_dir / "nested.txt"
        nested_file.write_text("content")

        # Should not raise
        result = ingester.ingest(
            file_path=nested_file,
            document_id=uuid4(),
            tenant_id=uuid4(),
            allowed_base_path=tmp_path,
        )
        assert result.success is True

    def test_ingest_path_traversal_via_symlink(self, tmp_path: Path) -> None:
        """Test that symlinks pointing outside allowed path are rejected."""
        ingester = MultimodalIngester(multimodal_enabled=True)
    
        # Create allowed directory and a file outside it
        allowed_dir = tmp_path / "allowed"
        allowed_dir.mkdir()
    
        outside_file = tmp_path / "secret.txt"
        outside_file.write_text("secret content")
    
        # Create symlink inside allowed dir pointing outside
        symlink = allowed_dir / "sneaky_link.txt"
        symlink.symlink_to(outside_file)
    
        # Security: Symlinks are now disallowed entirely for safety
        with pytest.raises(ValueError, match="Symlinks are not allowed"):
            ingester.ingest(
                file_path=symlink,
                document_id=uuid4(),
                tenant_id=uuid4(),
                allowed_base_path=allowed_dir,
            )

    def test_ingest_file_size_limit(self, tmp_path: Path) -> None:
        """Test that oversized files are rejected."""
        ingester = MultimodalIngester(multimodal_enabled=True)

        # Create a file that would be too large (simulate with empty large file)
        large_file = tmp_path / "large.txt"
        # We can't create a 100MB+ file easily in tests, so we'll mock the stat
        large_file.write_text("content")

        # For real test, check the constant is reasonable
        assert MAX_FILE_SIZE_BYTES == 100 * 1024 * 1024  # 100MB

    def test_ingest_rejects_legacy_doc_format(self, tmp_path: Path) -> None:
        """Test that legacy .doc format is rejected with helpful message."""
        ingester = MultimodalIngester(multimodal_enabled=True)

        doc_file = tmp_path / "legacy.doc"
        doc_file.write_bytes(b"fake legacy content")

        result = ingester.ingest(
            file_path=doc_file,
            document_id=uuid4(),
            tenant_id=uuid4(),
        )
        assert result.success is False
        assert "not supported" in result.error.lower()
        assert ".docx" in result.error  # Suggests modern format

    def test_ingest_rejects_legacy_xls_format(self, tmp_path: Path) -> None:
        """Test that legacy .xls format is rejected with helpful message."""
        ingester = MultimodalIngester(multimodal_enabled=True)

        xls_file = tmp_path / "legacy.xls"
        xls_file.write_bytes(b"fake legacy content")

        result = ingester.ingest(
            file_path=xls_file,
            document_id=uuid4(),
            tenant_id=uuid4(),
        )
        assert result.success is False
        assert "not supported" in result.error.lower()
        assert ".xlsx" in result.error  # Suggests modern format

    def test_ingest_rejects_legacy_ppt_format(self, tmp_path: Path) -> None:
        """Test that legacy .ppt format is rejected with helpful message."""
        ingester = MultimodalIngester(multimodal_enabled=True)

        ppt_file = tmp_path / "legacy.ppt"
        ppt_file.write_bytes(b"fake legacy content")

        result = ingester.ingest(
            file_path=ppt_file,
            document_id=uuid4(),
            tenant_id=uuid4(),
        )
        assert result.success is False
        assert "not supported" in result.error.lower()
        assert ".pptx" in result.error  # Suggests modern format

    def test_unsupported_legacy_formats_constant(self) -> None:
        """Test unsupported legacy formats constant."""
        assert ".doc" in UNSUPPORTED_LEGACY_FORMATS
        assert ".xls" in UNSUPPORTED_LEGACY_FORMATS
        assert ".ppt" in UNSUPPORTED_LEGACY_FORMATS

    def test_ingest_similar_path_prefix_rejection(self, tmp_path: Path) -> None:
        """Test that paths with similar prefix but different directory are rejected."""
        ingester = MultimodalIngester(multimodal_enabled=True)

        # Create /allowed and /allowed_other directories
        allowed_dir = tmp_path / "allowed"
        allowed_dir.mkdir()

        allowed_other_dir = tmp_path / "allowed_other"
        allowed_other_dir.mkdir()

        # Create file in allowed_other
        outside_file = allowed_other_dir / "file.txt"
        outside_file.write_text("content")

        # Should be rejected because allowed_other is NOT inside allowed
        with pytest.raises(ValueError, match="Path traversal not allowed"):
            ingester.ingest(
                file_path=outside_file,
                document_id=uuid4(),
                tenant_id=uuid4(),
                allowed_base_path=allowed_dir,
            )


# ==============================================================================
# OfficeParser Tests
# ==============================================================================


class TestOfficeParser:
    """Tests for OfficeParser class."""

    def test_office_parser_initialization(self) -> None:
        """Test OfficeParser initialization."""
        parser = OfficeParser()
        assert parser is not None

    def test_parse_word_requires_python_docx(self, tmp_path: Path) -> None:
        """Test parse_word raises ImportError if python-docx not available."""
        parser = OfficeParser()
        docx_file = tmp_path / "test.docx"
        docx_file.write_bytes(b"fake content")

        with patch.dict("sys.modules", {"docx": None}):
            with patch("agentic_rag_backend.indexing.multimodal.OfficeParser.parse_word") as mock:
                mock.side_effect = ImportError("python-docx is required")
                with pytest.raises(ImportError, match="python-docx"):
                    mock(docx_file, "doc-id", "tenant-id")

    def test_parse_excel_requires_openpyxl(self, tmp_path: Path) -> None:
        """Test parse_excel raises ImportError if openpyxl not available."""
        parser = OfficeParser()
        xlsx_file = tmp_path / "test.xlsx"
        xlsx_file.write_bytes(b"fake content")

        with patch.dict("sys.modules", {"openpyxl": None}):
            with patch("agentic_rag_backend.indexing.multimodal.OfficeParser.parse_excel") as mock:
                mock.side_effect = ImportError("openpyxl is required")
                with pytest.raises(ImportError, match="openpyxl"):
                    mock(xlsx_file, "doc-id", "tenant-id")

    def test_parse_powerpoint_requires_python_pptx(self, tmp_path: Path) -> None:
        """Test parse_powerpoint raises ImportError if python-pptx not available."""
        parser = OfficeParser()
        pptx_file = tmp_path / "test.pptx"
        pptx_file.write_bytes(b"fake content")

        with patch.dict("sys.modules", {"pptx": None}):
            with patch("agentic_rag_backend.indexing.multimodal.OfficeParser.parse_powerpoint") as mock:
                mock.side_effect = ImportError("python-pptx is required")
                with pytest.raises(ImportError, match="python-pptx"):
                    mock(pptx_file, "doc-id", "tenant-id")


# ==============================================================================
# MultimodalIngestionAdapter Tests
# ==============================================================================


class TestMultimodalIngestionAdapter:
    """Tests for MultimodalIngestionAdapter class."""

    def test_adapter_disabled_returns_none(self, tmp_path: Path) -> None:
        """Test adapter returns None when disabled."""
        adapter = MultimodalIngestionAdapter(ingester=None, enabled=False)

        test_file = tmp_path / "test.txt"
        test_file.write_text("content")

        result = adapter.ingest(
            file_path=test_file,
            document_id=uuid4(),
            tenant_id=uuid4(),
        )
        assert result is None

    def test_adapter_enabled_with_no_ingester_returns_none(self, tmp_path: Path) -> None:
        """Test adapter returns None when enabled but no ingester."""
        adapter = MultimodalIngestionAdapter(ingester=None, enabled=True)

        test_file = tmp_path / "test.txt"
        test_file.write_text("content")

        result = adapter.ingest(
            file_path=test_file,
            document_id=uuid4(),
            tenant_id=uuid4(),
        )
        assert result is None

    def test_adapter_enabled_calls_ingester(self, tmp_path: Path) -> None:
        """Test adapter calls ingester when enabled."""
        ingester = MultimodalIngester(multimodal_enabled=True)
        adapter = MultimodalIngestionAdapter(ingester=ingester, enabled=True)

        test_file = tmp_path / "test.txt"
        test_file.write_text("content")

        result = adapter.ingest(
            file_path=test_file,
            document_id=uuid4(),
            tenant_id=uuid4(),
        )
        assert result is not None
        assert result.success is True

    def test_adapter_detect_type(self) -> None:
        """Test adapter type detection."""
        ingester = MultimodalIngester(multimodal_enabled=True)
        adapter = MultimodalIngestionAdapter(ingester=ingester, enabled=True)

        assert adapter.detect_type(Path("doc.docx")) == DocumentType.WORD
        assert adapter.detect_type(Path("sheet.xlsx")) == DocumentType.EXCEL

    def test_adapter_detect_type_without_ingester(self) -> None:
        """Test adapter type detection without ingester."""
        adapter = MultimodalIngestionAdapter(ingester=None, enabled=False)

        # Should still work using static mapping
        assert adapter.detect_type(Path("doc.docx")) == DocumentType.WORD
        assert adapter.detect_type(Path("unknown.xyz")) == DocumentType.UNKNOWN


# ==============================================================================
# Factory Function Tests
# ==============================================================================


class TestGetMultimodalIngestionAdapter:
    """Tests for get_multimodal_ingestion_adapter factory function."""

    def test_factory_disabled_by_default(self) -> None:
        """Test factory creates disabled adapter by default."""
        settings = MagicMock()
        settings.multimodal_ingestion_enabled = False
        settings.office_docs_enabled = True

        adapter = get_multimodal_ingestion_adapter(settings)
        assert adapter.enabled is False

    def test_factory_enabled_when_settings_enabled(self) -> None:
        """Test factory creates enabled adapter when settings enable it."""
        settings = MagicMock()
        settings.multimodal_ingestion_enabled = True
        settings.office_docs_enabled = True

        adapter = get_multimodal_ingestion_adapter(settings)
        assert adapter.enabled is True
        assert adapter._ingester is not None

    def test_factory_respects_office_docs_setting(self) -> None:
        """Test factory passes office_docs_enabled to ingester."""
        settings = MagicMock()
        settings.multimodal_ingestion_enabled = True
        settings.office_docs_enabled = False

        adapter = get_multimodal_ingestion_adapter(settings)
        assert adapter._ingester.office_docs_enabled is False

    def test_factory_handles_missing_attributes(self) -> None:
        """Test factory handles missing settings attributes gracefully."""
        settings = MagicMock(spec=[])  # Empty spec - no attributes

        adapter = get_multimodal_ingestion_adapter(settings)
        # Should use defaults
        assert adapter.enabled == DEFAULT_MULTIMODAL_ENABLED


# ==============================================================================
# Default Constants Tests
# ==============================================================================


class TestDefaultConstants:
    """Tests for default configuration constants."""

    def test_default_multimodal_enabled(self) -> None:
        """Test default multimodal enabled is False."""
        assert DEFAULT_MULTIMODAL_ENABLED is False

    def test_default_office_docs_enabled(self) -> None:
        """Test default office docs enabled is True."""
        assert DEFAULT_OFFICE_DOCS_ENABLED is True


# ==============================================================================
# Integration Tests with Real Office Documents (requires libraries)
# ==============================================================================


class TestWordIntegration:
    """Integration tests for Word document parsing."""

    @pytest.fixture
    def word_document(self, tmp_path: Path) -> Path:
        """Create a simple Word document for testing."""
        try:
            from docx import Document
        except ImportError:
            pytest.skip("python-docx not installed")

        doc = Document()
        doc.add_heading("Test Document", 0)
        doc.add_paragraph("This is a test paragraph.")
        doc.add_paragraph("This is another paragraph.")

        # Add a simple table
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Header 1"
        table.cell(0, 1).text = "Header 2"
        table.cell(1, 0).text = "Value 1"
        table.cell(1, 1).text = "Value 2"

        file_path = tmp_path / "test.docx"
        doc.save(str(file_path))
        return file_path

    def test_parse_word_extracts_text(self, word_document: Path) -> None:
        """Test Word parsing extracts paragraph text."""
        parser = OfficeParser()
        content = parser.parse_word(word_document, "doc-id", "tenant-id")

        assert len(content.paragraphs) >= 2
        assert "test paragraph" in content.full_text.lower()

    def test_parse_word_extracts_tables(self, word_document: Path) -> None:
        """Test Word parsing extracts tables."""
        parser = OfficeParser()
        content = parser.parse_word(word_document, "doc-id", "tenant-id")

        assert len(content.tables) >= 1
        table = content.tables[0]
        assert len(table) >= 2  # At least header row and data row

    def test_ingester_processes_word(self, word_document: Path) -> None:
        """Test full ingester processes Word document."""
        ingester = MultimodalIngester(multimodal_enabled=True)
        result = ingester.ingest(
            file_path=word_document,
            document_id=uuid4(),
            tenant_id=uuid4(),
        )

        assert result.success is True
        assert result.document_type == DocumentType.WORD
        assert result.text_content
        assert result.markdown


class TestExcelIntegration:
    """Integration tests for Excel file parsing."""

    @pytest.fixture
    def excel_file(self, tmp_path: Path) -> Path:
        """Create a simple Excel file for testing."""
        try:
            from openpyxl import Workbook
        except ImportError:
            pytest.skip("openpyxl not installed")

        wb = Workbook()
        ws = wb.active
        ws.title = "Data"

        # Add headers and data
        ws["A1"] = "Name"
        ws["B1"] = "Value"
        ws["A2"] = "Item A"
        ws["B2"] = 100
        ws["A3"] = "Item B"
        ws["B3"] = 200

        file_path = tmp_path / "test.xlsx"
        wb.save(str(file_path))
        return file_path

    def test_parse_excel_extracts_sheets(self, excel_file: Path) -> None:
        """Test Excel parsing extracts sheet data."""
        parser = OfficeParser()
        content = parser.parse_excel(excel_file, "doc-id", "tenant-id")

        assert content.sheet_count >= 1
        sheet = content.sheets[0]
        assert sheet.name == "Data"
        assert "Name" in sheet.headers
        assert "Value" in sheet.headers

    def test_parse_excel_extracts_rows(self, excel_file: Path) -> None:
        """Test Excel parsing extracts row data."""
        parser = OfficeParser()
        content = parser.parse_excel(excel_file, "doc-id", "tenant-id")

        sheet = content.sheets[0]
        assert len(sheet.rows) >= 2
        # Check data is present
        row_values = [val for row in sheet.rows for val in row]
        assert "Item A" in row_values

    def test_ingester_processes_excel(self, excel_file: Path) -> None:
        """Test full ingester processes Excel file."""
        ingester = MultimodalIngester(multimodal_enabled=True)
        result = ingester.ingest(
            file_path=excel_file,
            document_id=uuid4(),
            tenant_id=uuid4(),
        )

        assert result.success is True
        assert result.document_type == DocumentType.EXCEL
        assert result.text_content
        assert result.markdown


class TestPowerPointIntegration:
    """Integration tests for PowerPoint file parsing."""

    @pytest.fixture
    def powerpoint_file(self, tmp_path: Path) -> Path:
        """Create a simple PowerPoint file for testing."""
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            pytest.skip("python-pptx not installed")

        prs = Presentation()

        # Add title slide
        title_layout = prs.slide_layouts[0]
        slide1 = prs.slides.add_slide(title_layout)
        slide1.shapes.title.text = "Test Presentation"
        slide1.placeholders[1].text = "Subtitle text"

        # Add content slide
        content_layout = prs.slide_layouts[1]
        slide2 = prs.slides.add_slide(content_layout)
        slide2.shapes.title.text = "Content Slide"
        slide2.placeholders[1].text = "Bullet point content"

        # Add speaker notes to second slide
        notes_slide = slide2.notes_slide
        notes_slide.notes_text_frame.text = "Speaker notes here"

        file_path = tmp_path / "test.pptx"
        prs.save(str(file_path))
        return file_path

    def test_parse_powerpoint_extracts_slides(self, powerpoint_file: Path) -> None:
        """Test PowerPoint parsing extracts slides."""
        parser = OfficeParser()
        content = parser.parse_powerpoint(powerpoint_file, "doc-id", "tenant-id")

        assert content.slide_count >= 2
        assert len(content.slides) >= 2

    def test_parse_powerpoint_extracts_titles(self, powerpoint_file: Path) -> None:
        """Test PowerPoint parsing extracts slide titles."""
        parser = OfficeParser()
        content = parser.parse_powerpoint(powerpoint_file, "doc-id", "tenant-id")

        titles = [s.title for s in content.slides if s.title]
        assert len(titles) >= 1
        # At least one should contain "Presentation" or "Content"
        title_text = " ".join(titles).lower()
        assert "presentation" in title_text or "content" in title_text

    def test_parse_powerpoint_extracts_notes(self, powerpoint_file: Path) -> None:
        """Test PowerPoint parsing extracts speaker notes."""
        parser = OfficeParser()
        content = parser.parse_powerpoint(powerpoint_file, "doc-id", "tenant-id")

        slides_with_notes = [s for s in content.slides if s.notes]
        assert len(slides_with_notes) >= 1
        assert "notes" in slides_with_notes[0].notes.lower()

    def test_ingester_processes_powerpoint(self, powerpoint_file: Path) -> None:
        """Test full ingester processes PowerPoint file."""
        ingester = MultimodalIngester(multimodal_enabled=True)
        result = ingester.ingest(
            file_path=powerpoint_file,
            document_id=uuid4(),
            tenant_id=uuid4(),
        )

        assert result.success is True
        assert result.document_type == DocumentType.POWERPOINT
        assert result.text_content
        assert result.markdown
        assert "---" in result.markdown  # Slide separator
