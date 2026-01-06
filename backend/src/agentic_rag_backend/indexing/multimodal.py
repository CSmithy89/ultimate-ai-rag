"""Multimodal Document Ingestion (Story 20-D2).

This module provides multimodal document ingestion capabilities for:
- Word documents (.docx) - text and tables extraction
- Excel files (.xlsx) - all sheets with cell data
- PowerPoint files (.pptx) - slides, text, and notes
- Automatic document type detection

Key Features:
- DocumentType enum for supported formats
- OfficeParser class for Microsoft Office documents
- MultimodalIngester for unified ingestion entry point
- Feature flags: MULTIMODAL_INGESTION_ENABLED, OFFICE_DOCS_ENABLED

Configuration:
- MULTIMODAL_INGESTION_ENABLED: Enable/disable multimodal processing (default: false)
- OFFICE_DOCS_ENABLED: Enable Office document processing (default: true when multimodal enabled)

Performance target: <500ms latency for typical documents
"""

import hashlib
import mimetypes
import time
from dataclasses import dataclass, field
from enum import Enum
from pathlib import Path
from typing import Any, Optional
from uuid import UUID

import structlog

logger = structlog.get_logger(__name__)

# Default configuration
DEFAULT_MULTIMODAL_ENABLED = False
DEFAULT_OFFICE_DOCS_ENABLED = True

# Document size limits to prevent DoS attacks
MAX_FILE_SIZE_BYTES = 100 * 1024 * 1024  # 100MB default
MAX_EXCEL_ROWS = 100000
MAX_EXCEL_COLS = 1000
MAX_EXCEL_SHEETS = 100
MAX_POWERPOINT_SLIDES = 500
MAX_WORD_PARAGRAPHS = 50000

# Legacy Office formats not supported by modern parsing libraries
UNSUPPORTED_LEGACY_FORMATS = {".doc", ".xls", ".ppt"}


class DocumentType(Enum):
    """Supported document types for multimodal ingestion.

    Attributes:
        PDF: PDF documents (.pdf)
        WORD: Microsoft Word documents (.docx)
        EXCEL: Microsoft Excel spreadsheets (.xlsx)
        POWERPOINT: Microsoft PowerPoint presentations (.pptx)
        IMAGE: Image files (.png, .jpg, .jpeg, .gif, .webp)
        MARKDOWN: Markdown documents (.md)
        TEXT: Plain text files (.txt)
        UNKNOWN: Unknown or unsupported document type
    """

    PDF = "pdf"
    WORD = "word"
    EXCEL = "excel"
    POWERPOINT = "powerpoint"
    IMAGE = "image"
    MARKDOWN = "markdown"
    TEXT = "text"
    UNKNOWN = "unknown"


# File extension to DocumentType mapping
EXTENSION_TYPE_MAP: dict[str, DocumentType] = {
    ".pdf": DocumentType.PDF,
    ".docx": DocumentType.WORD,
    ".doc": DocumentType.WORD,  # Legacy Word format
    ".xlsx": DocumentType.EXCEL,
    ".xls": DocumentType.EXCEL,  # Legacy Excel format
    ".pptx": DocumentType.POWERPOINT,
    ".ppt": DocumentType.POWERPOINT,  # Legacy PowerPoint format
    ".png": DocumentType.IMAGE,
    ".jpg": DocumentType.IMAGE,
    ".jpeg": DocumentType.IMAGE,
    ".gif": DocumentType.IMAGE,
    ".webp": DocumentType.IMAGE,
    ".md": DocumentType.MARKDOWN,
    ".markdown": DocumentType.MARKDOWN,
    ".txt": DocumentType.TEXT,
}

# MIME type to DocumentType mapping
MIME_TYPE_MAP: dict[str, DocumentType] = {
    "application/pdf": DocumentType.PDF,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": DocumentType.WORD,
    "application/msword": DocumentType.WORD,
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": DocumentType.EXCEL,
    "application/vnd.ms-excel": DocumentType.EXCEL,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": DocumentType.POWERPOINT,
    "application/vnd.ms-powerpoint": DocumentType.POWERPOINT,
    "image/png": DocumentType.IMAGE,
    "image/jpeg": DocumentType.IMAGE,
    "image/gif": DocumentType.IMAGE,
    "image/webp": DocumentType.IMAGE,
    "text/markdown": DocumentType.MARKDOWN,
    "text/plain": DocumentType.TEXT,
}


@dataclass
class ExtractedCell:
    """Represents a cell from an Excel spreadsheet.

    Attributes:
        row: Row index (0-based)
        column: Column index (0-based)
        value: Cell value as string
        data_type: Cell data type (string, number, date, formula, etc.)
    """

    row: int
    column: int
    value: str
    data_type: str = "string"

    def to_dict(self) -> dict[str, Any]:
        """Serialize cell to dictionary."""
        return {
            "row": self.row,
            "column": self.column,
            "value": self.value,
            "data_type": self.data_type,
        }


@dataclass
class ExtractedSheet:
    """Represents an Excel spreadsheet sheet.

    Attributes:
        name: Sheet name
        index: Sheet index (0-based)
        headers: Column headers (from first row)
        rows: Data rows (excluding header row)
        cells: All cells with their positions
        row_count: Total number of rows
        column_count: Total number of columns
    """

    name: str
    index: int
    headers: list[str] = field(default_factory=list)
    rows: list[list[str]] = field(default_factory=list)
    cells: list[ExtractedCell] = field(default_factory=list)
    row_count: int = 0
    column_count: int = 0

    def to_dict(self) -> dict[str, Any]:
        """Serialize sheet to dictionary."""
        return {
            "name": self.name,
            "index": self.index,
            "headers": self.headers,
            "rows": self.rows,
            "cells": [c.to_dict() for c in self.cells],
            "row_count": self.row_count,
            "column_count": self.column_count,
        }

    def to_markdown(self) -> str:
        """Convert sheet to markdown table format."""
        if not self.headers and not self.rows:
            return ""

        lines = []
        lines.append(f"## {self.name}\n")

        if self.headers:
            lines.append("| " + " | ".join(self.headers) + " |")
            lines.append("|" + "|".join(["---"] * len(self.headers)) + "|")

        for row in self.rows:
            # Escape pipe characters
            escaped = [cell.replace("|", "\\|") for cell in row]
            lines.append("| " + " | ".join(escaped) + " |")

        return "\n".join(lines)


@dataclass
class ExtractedSlide:
    """Represents a PowerPoint slide.

    Attributes:
        number: Slide number (1-based)
        title: Slide title
        content: Main slide content/text
        notes: Speaker notes
        shapes_text: Text from all shapes
    """

    number: int
    title: Optional[str] = None
    content: str = ""
    notes: Optional[str] = None
    shapes_text: list[str] = field(default_factory=list)

    def to_dict(self) -> dict[str, Any]:
        """Serialize slide to dictionary."""
        return {
            "number": self.number,
            "title": self.title,
            "content": self.content,
            "notes": self.notes,
            "shapes_text": self.shapes_text,
        }

    def to_markdown(self) -> str:
        """Convert slide to markdown format."""
        lines = []
        lines.append(f"## Slide {self.number}")

        if self.title:
            lines.append(f"### {self.title}\n")

        if self.content:
            lines.append(self.content)

        if self.shapes_text:
            for text in self.shapes_text:
                if text.strip() and text != self.title and text != self.content:
                    lines.append(f"\n{text}")

        if self.notes:
            lines.append(f"\n**Notes:** {self.notes}")

        return "\n".join(lines)


@dataclass
class WordContent:
    """Represents content extracted from a Word document.

    Attributes:
        paragraphs: List of paragraph texts
        tables: Extracted tables as list of rows
        full_text: Complete document text
        metadata: Document metadata (author, title, etc.)
    """

    paragraphs: list[str] = field(default_factory=list)
    tables: list[list[list[str]]] = field(default_factory=list)
    full_text: str = ""
    metadata: dict[str, str] = field(default_factory=dict)

    def to_dict(self) -> dict[str, Any]:
        """Serialize Word content to dictionary."""
        return {
            "paragraphs": self.paragraphs,
            "tables": self.tables,
            "full_text": self.full_text,
            "metadata": self.metadata,
        }


@dataclass
class ExcelContent:
    """Represents content extracted from an Excel file.

    Attributes:
        sheets: List of extracted sheets
        sheet_count: Total number of sheets
    """

    sheets: list[ExtractedSheet] = field(default_factory=list)
    sheet_count: int = 0

    def to_dict(self) -> dict[str, Any]:
        """Serialize Excel content to dictionary."""
        return {
            "sheets": [s.to_dict() for s in self.sheets],
            "sheet_count": self.sheet_count,
        }


@dataclass
class PowerPointContent:
    """Represents content extracted from a PowerPoint file.

    Attributes:
        slides: List of extracted slides
        slide_count: Total number of slides
        full_text: Combined text from all slides
    """

    slides: list[ExtractedSlide] = field(default_factory=list)
    slide_count: int = 0
    full_text: str = ""

    def to_dict(self) -> dict[str, Any]:
        """Serialize PowerPoint content to dictionary."""
        return {
            "slides": [s.to_dict() for s in self.slides],
            "slide_count": self.slide_count,
            "full_text": self.full_text,
        }


@dataclass
class MultimodalIngestionResult:
    """Result of multimodal document ingestion.

    Attributes:
        document_id: Document identifier
        tenant_id: Tenant identifier
        document_type: Detected document type
        file_path: Source file path
        content: Extracted content (type-specific)
        text_content: Unified text representation
        markdown: Markdown representation
        processing_time_ms: Processing time in milliseconds
        success: Whether ingestion was successful
        error: Error message if failed
        metadata: Additional metadata
    """

    document_id: str
    tenant_id: str
    document_type: DocumentType
    file_path: str
    content: Optional[dict[str, Any]] = None
    text_content: str = ""
    markdown: str = ""
    processing_time_ms: int = 0
    success: bool = True
    error: Optional[str] = None
    metadata: dict[str, Any] = field(default_factory=dict)

    def to_dict(self) -> dict[str, Any]:
        """Serialize result to dictionary."""
        return {
            "document_id": self.document_id,
            "tenant_id": self.tenant_id,
            "document_type": self.document_type.value,
            "file_path": self.file_path,
            "content": self.content,
            "text_content": self.text_content,
            "markdown": self.markdown,
            "processing_time_ms": self.processing_time_ms,
            "success": self.success,
            "error": self.error,
            "metadata": self.metadata,
        }


class OfficeParser:
    """Parser for Microsoft Office documents.

    This class handles parsing of Word, Excel, and PowerPoint documents
    using python-docx, openpyxl, and python-pptx libraries.
    """

    def __init__(self) -> None:
        """Initialize OfficeParser."""
        logger.debug("office_parser_initialized")

    def parse_word(
        self,
        file_path: Path,
        document_id: str,
        tenant_id: str,
    ) -> WordContent:
        """Parse a Word document (.docx).

        Extracts text from paragraphs and tables.

        Args:
            file_path: Path to Word document
            document_id: Document identifier
            tenant_id: Tenant identifier

        Returns:
            WordContent with extracted content

        Raises:
            ImportError: If python-docx is not installed
            Exception: If parsing fails
        """
        try:
            from docx import Document
        except ImportError as e:
            logger.error("python_docx_not_installed", error=str(e))
            raise ImportError(
                "python-docx is required for Word document parsing. "
                "Install with: pip install python-docx"
            ) from e

        logger.info(
            "parse_word_started",
            file_path=str(file_path),
            document_id=document_id,
        )

        doc = Document(str(file_path))

        # Extract paragraphs with size limit
        paragraphs: list[str] = []
        for i, para in enumerate(doc.paragraphs):
            if i >= MAX_WORD_PARAGRAPHS:
                logger.warning(
                    "word_paragraph_limit_exceeded",
                    document_id=document_id,
                    limit=MAX_WORD_PARAGRAPHS,
                )
                break
            if para.text.strip():
                paragraphs.append(para.text)

        # Extract tables
        tables: list[list[list[str]]] = []
        for table in doc.tables:
            table_data: list[list[str]] = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            if table_data:
                tables.append(table_data)

        # Build full text
        full_text = "\n\n".join(paragraphs)

        # Extract metadata
        metadata: dict[str, str] = {}
        core_props = doc.core_properties
        if core_props.author:
            metadata["author"] = core_props.author
        if core_props.title:
            metadata["title"] = core_props.title
        if core_props.subject:
            metadata["subject"] = core_props.subject
        if core_props.created:
            metadata["created"] = str(core_props.created)
        if core_props.modified:
            metadata["modified"] = str(core_props.modified)

        logger.info(
            "parse_word_completed",
            document_id=document_id,
            paragraphs=len(paragraphs),
            tables=len(tables),
        )

        return WordContent(
            paragraphs=paragraphs,
            tables=tables,
            full_text=full_text,
            metadata=metadata,
        )

    def parse_excel(
        self,
        file_path: Path,
        document_id: str,
        tenant_id: str,
    ) -> ExcelContent:
        """Parse an Excel file (.xlsx).

        Extracts all sheets with cell data.

        Args:
            file_path: Path to Excel file
            document_id: Document identifier
            tenant_id: Tenant identifier

        Returns:
            ExcelContent with extracted sheets

        Raises:
            ImportError: If openpyxl is not installed
            Exception: If parsing fails
        """
        try:
            from openpyxl import load_workbook
        except ImportError as e:
            logger.error("openpyxl_not_installed", error=str(e))
            raise ImportError(
                "openpyxl is required for Excel file parsing. "
                "Install with: pip install openpyxl"
            ) from e

        logger.info(
            "parse_excel_started",
            file_path=str(file_path),
            document_id=document_id,
        )

        wb = load_workbook(str(file_path), data_only=True)

        try:
            sheets: list[ExtractedSheet] = []
            for sheet_idx, sheet_name in enumerate(wb.sheetnames):
                if sheet_idx >= MAX_EXCEL_SHEETS:
                    logger.warning(
                        "excel_sheet_limit_exceeded",
                        document_id=document_id,
                        limit=MAX_EXCEL_SHEETS,
                    )
                    break

                ws = wb[sheet_name]

                # Get dimensions
                max_row = min(ws.max_row or 0, MAX_EXCEL_ROWS)
                max_col = min(ws.max_column or 0, MAX_EXCEL_COLS)

                if max_row == 0 or max_col == 0:
                    continue

                # Extract data
                cells: list[ExtractedCell] = []
                all_rows: list[list[str]] = []

                for row_idx, row in enumerate(ws.iter_rows(max_row=max_row, max_col=max_col)):
                    if row_idx >= MAX_EXCEL_ROWS:
                        break
                    row_values: list[str] = []
                    for col_idx, cell in enumerate(row):
                        if col_idx >= MAX_EXCEL_COLS:
                            break
                        value = str(cell.value) if cell.value is not None else ""
                        data_type = type(cell.value).__name__ if cell.value is not None else "empty"
                        row_values.append(value)
                        cells.append(ExtractedCell(
                            row=row_idx,
                            column=col_idx,
                            value=value,
                            data_type=data_type,
                        ))
                    all_rows.append(row_values)

                # First row as headers
                headers = all_rows[0] if all_rows else []
                data_rows = all_rows[1:] if len(all_rows) > 1 else []

                sheets.append(ExtractedSheet(
                    name=sheet_name,
                    index=sheet_idx,
                    headers=headers,
                    rows=data_rows,
                    cells=cells,
                    row_count=len(all_rows),
                    column_count=max_col,
                ))

            logger.info(
                "parse_excel_completed",
                document_id=document_id,
                sheets=len(sheets),
            )

            return ExcelContent(
                sheets=sheets,
                sheet_count=len(sheets),
            )
        finally:
            wb.close()

    def parse_powerpoint(
        self,
        file_path: Path,
        document_id: str,
        tenant_id: str,
    ) -> PowerPointContent:
        """Parse a PowerPoint file (.pptx).

        Extracts slides, text, and speaker notes.

        Args:
            file_path: Path to PowerPoint file
            document_id: Document identifier
            tenant_id: Tenant identifier

        Returns:
            PowerPointContent with extracted slides

        Raises:
            ImportError: If python-pptx is not installed
            Exception: If parsing fails
        """
        try:
            from pptx import Presentation
        except ImportError as e:
            logger.error("python_pptx_not_installed", error=str(e))
            raise ImportError(
                "python-pptx is required for PowerPoint file parsing. "
                "Install with: pip install python-pptx"
            ) from e

        logger.info(
            "parse_powerpoint_started",
            file_path=str(file_path),
            document_id=document_id,
        )

        prs = Presentation(str(file_path))

        slides: list[ExtractedSlide] = []
        all_text: list[str] = []

        for slide_idx, slide in enumerate(prs.slides):
            if slide_idx >= MAX_POWERPOINT_SLIDES:
                logger.warning(
                    "powerpoint_slide_limit_exceeded",
                    document_id=document_id,
                    limit=MAX_POWERPOINT_SLIDES,
                )
                break

            slide_number = slide_idx + 1
            title: Optional[str] = None
            content_parts: list[str] = []
            shapes_text: list[str] = []

            # Extract title and text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    shapes_text.append(text)

                    # Check if this is the title
                    if shape.has_text_frame:
                        if hasattr(shape, "is_title") and shape.is_title:
                            title = text
                        elif title is None and slide_idx == 0:
                            # First text on first slide might be title
                            title = text
                        else:
                            content_parts.append(text)

            content = "\n".join(content_parts)

            # Extract speaker notes
            notes: Optional[str] = None
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    notes = notes_text

            extracted_slide = ExtractedSlide(
                number=slide_number,
                title=title,
                content=content,
                notes=notes,
                shapes_text=shapes_text,
            )
            slides.append(extracted_slide)

            # Collect text for full_text
            if title:
                all_text.append(f"Slide {slide_number}: {title}")
            if content:
                all_text.append(content)

        full_text = "\n\n".join(all_text)

        logger.info(
            "parse_powerpoint_completed",
            document_id=document_id,
            slides=len(slides),
        )

        return PowerPointContent(
            slides=slides,
            slide_count=len(slides),
            full_text=full_text,
        )


class MultimodalIngester:
    """Unified multimodal document ingestion.

    This class provides the main entry point for ingesting various
    document types with automatic type detection and routing.
    """

    def __init__(
        self,
        office_parser: Optional[OfficeParser] = None,
        multimodal_enabled: bool = DEFAULT_MULTIMODAL_ENABLED,
        office_docs_enabled: bool = DEFAULT_OFFICE_DOCS_ENABLED,
    ) -> None:
        """Initialize MultimodalIngester.

        Args:
            office_parser: OfficeParser instance (created if not provided)
            multimodal_enabled: Whether multimodal ingestion is enabled
            office_docs_enabled: Whether Office document processing is enabled
        """
        self._office_parser = office_parser or OfficeParser()
        self.multimodal_enabled = multimodal_enabled
        self.office_docs_enabled = office_docs_enabled

        logger.debug(
            "multimodal_ingester_initialized",
            multimodal_enabled=multimodal_enabled,
            office_docs_enabled=office_docs_enabled,
        )

    def ingest(
        self,
        file_path: Path,
        document_id: UUID,
        tenant_id: UUID,
        allowed_base_path: Optional[Path] = None,
        document_type: Optional[DocumentType] = None,
    ) -> MultimodalIngestionResult:
        """Ingest a document with automatic type detection.

        This is the main entry point for multimodal ingestion.
        Routes to appropriate parser based on detected document type.

        IMPORTANT: Callers MUST verify tenant authorization before calling.
        This module does not perform authorization checks.

        Args:
            file_path: Path to document file
            document_id: Document identifier
            tenant_id: Tenant identifier
            allowed_base_path: Optional base directory for path traversal protection
            document_type: Optional document type (auto-detected if not provided)

        Returns:
            MultimodalIngestionResult with extracted content

        Raises:
            ValueError: If tenant_id is invalid or file_path fails validation
            FileNotFoundError: If file does not exist
        """
        start_time = time.perf_counter()

        # Validate tenant_id - must be non-null for multi-tenancy
        if tenant_id is None:
            raise ValueError("tenant_id is required for multi-tenancy")

        # Validate file path exists and is a file
        if not file_path.exists():
            raise FileNotFoundError(f"Document not found: {file_path}")
        if not file_path.is_file():
            raise ValueError(f"Path is not a file: {file_path}")

        # Path traversal protection - use relative_to() for proper containment check
        if allowed_base_path is not None:
            resolved_path = file_path.resolve()
            resolved_base = allowed_base_path.resolve()
            try:
                resolved_path.relative_to(resolved_base)
            except ValueError:
                logger.warning(
                    "path_traversal_attempt",
                    file_path=str(file_path),
                    allowed_base=str(allowed_base_path),
                    tenant_id=str(tenant_id),
                )
                raise ValueError(f"Path traversal not allowed: {file_path}")

        # File size validation to prevent DoS
        file_size = file_path.stat().st_size
        if file_size > MAX_FILE_SIZE_BYTES:
            raise ValueError(
                f"File size ({file_size} bytes) exceeds maximum allowed "
                f"({MAX_FILE_SIZE_BYTES} bytes)"
            )

        # Check if multimodal ingestion is enabled
        if not self.multimodal_enabled:
            logger.debug(
                "multimodal_ingestion_disabled",
                document_id=str(document_id),
            )
            processing_time_ms = int((time.perf_counter() - start_time) * 1000)
            return MultimodalIngestionResult(
                document_id=str(document_id),
                tenant_id=str(tenant_id),
                document_type=DocumentType.UNKNOWN,
                file_path=str(file_path),
                success=False,
                error="Multimodal ingestion is disabled",
                processing_time_ms=processing_time_ms,
            )

        # Detect document type if not provided
        detected_type = document_type or self._detect_type(file_path)

        logger.info(
            "multimodal_ingest_started",
            file_path=str(file_path),
            document_id=str(document_id),
            tenant_id=str(tenant_id),
            document_type=detected_type.value,
        )

        try:
            result = self._process_document(
                file_path=file_path,
                document_id=document_id,
                tenant_id=tenant_id,
                document_type=detected_type,
            )

            processing_time_ms = int((time.perf_counter() - start_time) * 1000)
            result.processing_time_ms = processing_time_ms

            logger.info(
                "multimodal_ingest_completed",
                document_id=str(document_id),
                document_type=detected_type.value,
                processing_time_ms=processing_time_ms,
                success=result.success,
            )

            return result

        except Exception as e:
            processing_time_ms = int((time.perf_counter() - start_time) * 1000)
            logger.error(
                "multimodal_ingest_failed",
                document_id=str(document_id),
                document_type=detected_type.value,
                error=str(e),
            )
            return MultimodalIngestionResult(
                document_id=str(document_id),
                tenant_id=str(tenant_id),
                document_type=detected_type,
                file_path=str(file_path),
                success=False,
                error=str(e),
                processing_time_ms=processing_time_ms,
            )

    def _detect_type(self, file_path: Path) -> DocumentType:
        """Auto-detect document type from file extension or MIME type.

        Args:
            file_path: Path to document file

        Returns:
            Detected DocumentType
        """
        # Try extension first
        extension = file_path.suffix.lower()
        if extension in EXTENSION_TYPE_MAP:
            return EXTENSION_TYPE_MAP[extension]

        # Try MIME type
        mime_type, _ = mimetypes.guess_type(str(file_path))
        if mime_type and mime_type in MIME_TYPE_MAP:
            return MIME_TYPE_MAP[mime_type]

        return DocumentType.UNKNOWN

    def _process_document(
        self,
        file_path: Path,
        document_id: UUID,
        tenant_id: UUID,
        document_type: DocumentType,
    ) -> MultimodalIngestionResult:
        """Process document based on its type.

        Args:
            file_path: Path to document
            document_id: Document identifier
            tenant_id: Tenant identifier
            document_type: Detected document type

        Returns:
            MultimodalIngestionResult with processed content
        """
        doc_id_str = str(document_id)
        tenant_id_str = str(tenant_id)

        # Check for unsupported legacy Office formats
        extension = file_path.suffix.lower()
        if extension in UNSUPPORTED_LEGACY_FORMATS:
            return MultimodalIngestionResult(
                document_id=doc_id_str,
                tenant_id=tenant_id_str,
                document_type=document_type,
                file_path=str(file_path),
                success=False,
                error=f"Legacy Office format {extension} is not supported. "
                      f"Please convert to the modern format (.docx, .xlsx, .pptx).",
            )

        # Handle Office documents
        if document_type in (DocumentType.WORD, DocumentType.EXCEL, DocumentType.POWERPOINT):
            if not self.office_docs_enabled:
                return MultimodalIngestionResult(
                    document_id=doc_id_str,
                    tenant_id=tenant_id_str,
                    document_type=document_type,
                    file_path=str(file_path),
                    success=False,
                    error="Office document processing is disabled",
                )

        if document_type == DocumentType.WORD:
            return self._process_word(file_path, doc_id_str, tenant_id_str)
        elif document_type == DocumentType.EXCEL:
            return self._process_excel(file_path, doc_id_str, tenant_id_str)
        elif document_type == DocumentType.POWERPOINT:
            return self._process_powerpoint(file_path, doc_id_str, tenant_id_str)
        elif document_type == DocumentType.MARKDOWN:
            return self._process_markdown(file_path, doc_id_str, tenant_id_str)
        elif document_type == DocumentType.TEXT:
            return self._process_text(file_path, doc_id_str, tenant_id_str)
        elif document_type == DocumentType.PDF:
            # PDF processing should use EnhancedDoclingParser
            return MultimodalIngestionResult(
                document_id=doc_id_str,
                tenant_id=tenant_id_str,
                document_type=document_type,
                file_path=str(file_path),
                success=False,
                error="PDF processing should use EnhancedDoclingParser",
                metadata={"hint": "Use enhanced_docling.py for PDF files"},
            )
        elif document_type == DocumentType.IMAGE:
            # Image processing not implemented in this story
            return MultimodalIngestionResult(
                document_id=doc_id_str,
                tenant_id=tenant_id_str,
                document_type=document_type,
                file_path=str(file_path),
                success=False,
                error="Image processing not yet implemented",
                metadata={"hint": "IMAGE_INGESTION_ENABLED is reserved for future use"},
            )
        else:
            return MultimodalIngestionResult(
                document_id=doc_id_str,
                tenant_id=tenant_id_str,
                document_type=document_type,
                file_path=str(file_path),
                success=False,
                error=f"Unsupported document type: {document_type.value}",
            )

    def _process_word(
        self,
        file_path: Path,
        document_id: str,
        tenant_id: str,
    ) -> MultimodalIngestionResult:
        """Process Word document."""
        content = self._office_parser.parse_word(file_path, document_id, tenant_id)

        # Generate markdown representation
        markdown_parts = []
        for para in content.paragraphs:
            markdown_parts.append(para)

        for i, table in enumerate(content.tables):
            markdown_parts.append(f"\n### Table {i + 1}\n")
            if table:
                # First row as headers
                headers = table[0]
                markdown_parts.append("| " + " | ".join(headers) + " |")
                markdown_parts.append("|" + "|".join(["---"] * len(headers)) + "|")
                for row in table[1:]:
                    escaped = [cell.replace("|", "\\|") for cell in row]
                    markdown_parts.append("| " + " | ".join(escaped) + " |")

        markdown = "\n\n".join(markdown_parts)

        return MultimodalIngestionResult(
            document_id=document_id,
            tenant_id=tenant_id,
            document_type=DocumentType.WORD,
            file_path=str(file_path),
            content=content.to_dict(),
            text_content=content.full_text,
            markdown=markdown,
            success=True,
            metadata={
                "paragraph_count": len(content.paragraphs),
                "table_count": len(content.tables),
                **content.metadata,
            },
        )

    def _process_excel(
        self,
        file_path: Path,
        document_id: str,
        tenant_id: str,
    ) -> MultimodalIngestionResult:
        """Process Excel file."""
        content = self._office_parser.parse_excel(file_path, document_id, tenant_id)

        # Generate markdown representation
        markdown_parts = []
        text_parts = []

        for sheet in content.sheets:
            markdown_parts.append(sheet.to_markdown())

            # Plain text representation
            text_parts.append(f"Sheet: {sheet.name}")
            if sheet.headers:
                text_parts.append("\t".join(sheet.headers))
            for row in sheet.rows:
                text_parts.append("\t".join(row))

        markdown = "\n\n".join(markdown_parts)
        text_content = "\n".join(text_parts)

        return MultimodalIngestionResult(
            document_id=document_id,
            tenant_id=tenant_id,
            document_type=DocumentType.EXCEL,
            file_path=str(file_path),
            content=content.to_dict(),
            text_content=text_content,
            markdown=markdown,
            success=True,
            metadata={
                "sheet_count": content.sheet_count,
                "sheet_names": [s.name for s in content.sheets],
            },
        )

    def _process_powerpoint(
        self,
        file_path: Path,
        document_id: str,
        tenant_id: str,
    ) -> MultimodalIngestionResult:
        """Process PowerPoint file."""
        content = self._office_parser.parse_powerpoint(file_path, document_id, tenant_id)

        # Generate markdown representation
        markdown_parts = []
        for slide in content.slides:
            markdown_parts.append(slide.to_markdown())

        markdown = "\n\n---\n\n".join(markdown_parts)

        return MultimodalIngestionResult(
            document_id=document_id,
            tenant_id=tenant_id,
            document_type=DocumentType.POWERPOINT,
            file_path=str(file_path),
            content=content.to_dict(),
            text_content=content.full_text,
            markdown=markdown,
            success=True,
            metadata={
                "slide_count": content.slide_count,
            },
        )

    def _read_text_file(
        self,
        file_path: Path,
        document_id: str,
    ) -> tuple[str, Optional[str]]:
        """Read text file with encoding fallback.

        Attempts UTF-8 first, then falls back to common encodings.

        Args:
            file_path: Path to text file
            document_id: Document ID for logging

        Returns:
            Tuple of (content, encoding_used) or (None, error_message)
        """
        encodings = ["utf-8", "utf-8-sig", "latin-1", "cp1252"]

        for encoding in encodings:
            try:
                content = file_path.read_text(encoding=encoding)
                if encoding != "utf-8":
                    logger.info(
                        "text_file_encoding_fallback",
                        document_id=document_id,
                        encoding=encoding,
                    )
                return content, None
            except UnicodeDecodeError:
                continue

        return "", "Unable to decode file with supported encodings (utf-8, latin-1, cp1252)"

    def _process_markdown(
        self,
        file_path: Path,
        document_id: str,
        tenant_id: str,
    ) -> MultimodalIngestionResult:
        """Process Markdown file."""
        text_content, error = self._read_text_file(file_path, document_id)
        if error:
            return MultimodalIngestionResult(
                document_id=document_id,
                tenant_id=tenant_id,
                document_type=DocumentType.MARKDOWN,
                file_path=str(file_path),
                success=False,
                error=error,
            )

        return MultimodalIngestionResult(
            document_id=document_id,
            tenant_id=tenant_id,
            document_type=DocumentType.MARKDOWN,
            file_path=str(file_path),
            content={"raw_text": text_content},
            text_content=text_content,
            markdown=text_content,  # Already markdown
            success=True,
            metadata={
                "char_count": len(text_content),
            },
        )

    def _process_text(
        self,
        file_path: Path,
        document_id: str,
        tenant_id: str,
    ) -> MultimodalIngestionResult:
        """Process plain text file."""
        text_content, error = self._read_text_file(file_path, document_id)
        if error:
            return MultimodalIngestionResult(
                document_id=document_id,
                tenant_id=tenant_id,
                document_type=DocumentType.TEXT,
                file_path=str(file_path),
                success=False,
                error=error,
            )

        return MultimodalIngestionResult(
            document_id=document_id,
            tenant_id=tenant_id,
            document_type=DocumentType.TEXT,
            file_path=str(file_path),
            content={"raw_text": text_content},
            text_content=text_content,
            markdown=f"```\n{text_content}\n```",  # Wrap in code block
            success=True,
            metadata={
                "char_count": len(text_content),
            },
        )

    def _generate_id(self, id_string: str) -> str:
        """Generate deterministic ID from string.

        Args:
            id_string: String to hash

        Returns:
            Short hash-based ID
        """
        hash_digest = hashlib.sha256(id_string.encode()).hexdigest()[:16]
        return hash_digest


class MultimodalIngestionAdapter:
    """Adapter for MultimodalIngester with feature flag support.

    This adapter provides:
    - Feature flag checking (MULTIMODAL_INGESTION_ENABLED)
    - Configuration from Settings
    - Graceful fallback when disabled
    """

    def __init__(
        self,
        ingester: Optional[MultimodalIngester],
        enabled: bool = False,
    ) -> None:
        """Initialize MultimodalIngestionAdapter.

        Args:
            ingester: MultimodalIngester instance (or None if disabled)
            enabled: Whether multimodal ingestion is enabled
        """
        self._ingester = ingester
        self.enabled = enabled

    def ingest(
        self,
        file_path: Path,
        document_id: UUID,
        tenant_id: UUID,
        allowed_base_path: Optional[Path] = None,
        document_type: Optional[DocumentType] = None,
    ) -> Optional[MultimodalIngestionResult]:
        """Ingest document if multimodal ingestion is enabled.

        Args:
            file_path: Path to document file
            document_id: Document identifier
            tenant_id: Tenant identifier
            allowed_base_path: Optional base directory for path traversal protection
            document_type: Optional document type

        Returns:
            MultimodalIngestionResult if enabled and successful, None otherwise
        """
        if not self.enabled or not self._ingester:
            logger.debug(
                "multimodal_ingestion_adapter_disabled",
                enabled=self.enabled,
                has_ingester=self._ingester is not None,
            )
            return None

        return self._ingester.ingest(
            file_path=file_path,
            document_id=document_id,
            tenant_id=tenant_id,
            allowed_base_path=allowed_base_path,
            document_type=document_type,
        )

    def detect_type(self, file_path: Path) -> DocumentType:
        """Detect document type.

        Args:
            file_path: Path to document file

        Returns:
            Detected DocumentType
        """
        if not self._ingester:
            # Use static detection
            extension = file_path.suffix.lower()
            if extension in EXTENSION_TYPE_MAP:
                return EXTENSION_TYPE_MAP[extension]
            return DocumentType.UNKNOWN

        return self._ingester._detect_type(file_path)


def get_multimodal_ingestion_adapter(settings: Any) -> MultimodalIngestionAdapter:
    """Factory function to create MultimodalIngestionAdapter from settings.

    Args:
        settings: Application settings

    Returns:
        Configured MultimodalIngestionAdapter instance
    """
    enabled = getattr(settings, "multimodal_ingestion_enabled", DEFAULT_MULTIMODAL_ENABLED)
    office_docs_enabled = getattr(settings, "office_docs_enabled", DEFAULT_OFFICE_DOCS_ENABLED)

    ingester = None
    if enabled:
        ingester = MultimodalIngester(
            multimodal_enabled=enabled,
            office_docs_enabled=office_docs_enabled,
        )

    return MultimodalIngestionAdapter(
        ingester=ingester,
        enabled=enabled,
    )
